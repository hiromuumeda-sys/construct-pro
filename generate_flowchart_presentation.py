#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# プレゼンテーション作成
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# カラースキーム
COLOR_PURPLE = RGBColor(139, 71, 137)  # #8B4789
COLOR_PRIMARY = RGBColor(103, 80, 164)  # #6750A4
COLOR_START = RGBColor(76, 175, 80)  # 緑
COLOR_PROCESS = RGBColor(33, 150, 243)  # 青
COLOR_DECISION = RGBColor(255, 152, 0)  # オレンジ
COLOR_END = RGBColor(244, 67, 54)  # 赤
COLOR_DARK = RGBColor(49, 51, 56)
COLOR_LIGHT = RGBColor(245, 245, 245)

def add_title_slide(prs, title, subtitle):
    """タイトルスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PURPLE

    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    left = Inches(0.5)
    top = Inches(4.2)
    width = Inches(9)
    height = Inches(1)
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

def add_title_bar(slide, title):
    """スライドにタイトルバーを追加"""
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.8)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_PRIMARY
    title_shape.line.color.rgb = COLOR_PRIMARY

    title_frame = title_shape.text_frame
    title_frame.clear()
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(5)

def add_box(slide, left, top, width, height, text, color, text_color=RGBColor(255, 255, 255)):
    """四角形ボックスを追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    return shape

def add_oval(slide, left, top, width, height, text, color, text_color=RGBColor(255, 255, 255)):
    """楕円形を追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    return shape

def add_diamond(slide, left, top, width, height, text, color, text_color=RGBColor(255, 255, 255)):
    """ダイアモンド形（判断）を追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    return shape

def add_arrow(slide, x1, y1, x2, y2):
    """矢印を追加"""
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = COLOR_DARK
    connector.line.width = Pt(2)

# ============== スライド作成開始 ==============

# スライド 1: タイトル
add_title_slide(prs, "CONSTRUCT_PRO", "システムフロー & フローチャート")

# スライド 2: 新規案件登録フロー
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "新規案件登録フロー")

# 開始
add_oval(slide, Inches(3.75), Inches(1.2), Inches(1.5), Inches(0.5), "開始", COLOR_START)
add_arrow(slide, 5, 1.7, 5, 2.1)

# モーダル表示
add_box(slide, Inches(3.5), Inches(2.1), Inches(3), Inches(0.5), "案件登録モーダル表示", COLOR_PROCESS)
add_arrow(slide, 5, 2.6, 5, 3.0)

# 入力
add_box(slide, Inches(3.5), Inches(3.0), Inches(3), Inches(0.5), "案件情報を入力", COLOR_PROCESS)
add_arrow(slide, 5, 3.5, 5, 3.9)

# 登録ボタン
add_box(slide, Inches(3.5), Inches(3.9), Inches(3), Inches(0.5), "「登録」ボタンをクリック", COLOR_PROCESS)
add_arrow(slide, 5, 4.4, 5, 4.8)

# バリデーション
add_diamond(slide, Inches(3.75), Inches(4.8), Inches(2.5), Inches(0.6), "バリデーション\nOK?", COLOR_DECISION)
add_arrow(slide, 3.75, 5.1, 2.5, 5.3)
# エラーメッセージ（No の場合）
add_box(slide, Inches(1), Inches(5.2), Inches(2), Inches(0.4), "エラーメッセージ", COLOR_END)
add_arrow(slide, 2, 5.6, 2, 6.2)
add_arrow(slide, 2, 6.2, 5, 6.2)
add_arrow(slide, 5, 6.2, 5, 5.4)

# サーバー処理（OK の場合）
add_arrow(slide, 6.25, 5.1, 7.5, 5.3)
add_box(slide, Inches(6.5), Inches(5.2), Inches(2.5), Inches(0.4), "DB に INSERT", COLOR_PROCESS)
add_arrow(slide, 7.75, 5.6, 7.75, 6.0)

# キャッシュ再取得
add_box(slide, Inches(6.5), Inches(6.0), Inches(2.5), Inches(0.4), "キャッシュ再取得", COLOR_PROCESS)
add_arrow(slide, 7.75, 6.4, 5, 6.6)

# テーブル再描画
add_box(slide, Inches(3.5), Inches(6.6), Inches(3), Inches(0.4), "テーブル再描画", COLOR_PROCESS)
add_arrow(slide, 5, 7.0, 5, 7.0)

# 終了
add_oval(slide, Inches(3.75), Inches(6.9), Inches(1.5), Inches(0.5), "完了", COLOR_END)

# スライド 3: 注文登録フロー
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "注文登録フロー")

# 開始
add_oval(slide, Inches(3.75), Inches(1.2), Inches(1.5), Inches(0.5), "開始", COLOR_START)
add_arrow(slide, 5, 1.7, 5, 2.1)

# 案件選択
add_box(slide, Inches(3.5), Inches(2.1), Inches(3), Inches(0.5), "案件を選択", COLOR_PROCESS)
add_arrow(slide, 5, 2.6, 5, 3.0)

# 注文管理画面
add_box(slide, Inches(3.5), Inches(3.0), Inches(3), Inches(0.5), "注文管理画面を開く", COLOR_PROCESS)
add_arrow(slide, 5, 3.5, 5, 3.9)

# キャッシュ確認
add_diamond(slide, Inches(3.75), Inches(3.9), Inches(2.5), Inches(0.6), "キャッシュ\n準備?", COLOR_DECISION)

# キャッシュ初期化（No）
add_arrow(slide, 3.75, 4.5, 2.5, 4.7)
add_box(slide, Inches(0.8), Inches(4.6), Inches(2.4), Inches(0.4), "キャッシュ初期化", COLOR_PROCESS)
add_arrow(slide, 2, 5.0, 2, 5.4)
add_arrow(slide, 2, 5.4, 5, 5.6)

# 注文入力（Yes）
add_arrow(slide, 6.25, 4.2, 7.5, 4.4)
add_box(slide, Inches(6.5), Inches(4.3), Inches(2.5), Inches(0.5), "工事区分・外注先\n金額を入力", COLOR_PROCESS)
add_arrow(slide, 7.75, 4.8, 7.75, 5.2)

# 登録ボタン
add_box(slide, Inches(6.5), Inches(5.2), Inches(2.5), Inches(0.4), "「登録」をクリック", COLOR_PROCESS)
add_arrow(slide, 7.75, 5.6, 5, 5.8)

# テーブル再描画
add_box(slide, Inches(3.5), Inches(5.6), Inches(3), Inches(0.4), "テーブル再描画", COLOR_PROCESS)
add_arrow(slide, 5, 6.0, 5, 6.4)

# 終了
add_oval(slide, Inches(3.75), Inches(6.4), Inches(1.5), Inches(0.5), "完了", COLOR_END)

# スライド 4: 請求書発行フロー
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "請求書発行フロー")

# 開始
add_oval(slide, Inches(3.75), Inches(1.1), Inches(1.5), Inches(0.5), "開始", COLOR_START)
add_arrow(slide, 5, 1.6, 5, 2.0)

# ボタンクリック
add_box(slide, Inches(3.5), Inches(2.0), Inches(3), Inches(0.4), "「請求書を発行」クリック", COLOR_PROCESS)
add_arrow(slide, 5, 2.4, 5, 2.8)

# プロジェクトID取得
add_box(slide, Inches(3.5), Inches(2.8), Inches(3), Inches(0.4), "URL から project_id 取得", COLOR_PROCESS)
add_arrow(slide, 5, 3.2, 5, 3.6)

# API呼び出し
add_box(slide, Inches(3.5), Inches(3.6), Inches(3), Inches(0.5), "GET /api/invoice/project/:id\n呼び出し", COLOR_PROCESS)
add_arrow(slide, 5, 4.1, 5, 4.5)

# DB クエリ
add_box(slide, Inches(3.5), Inches(4.5), Inches(3), Inches(0.5), "案件・注文データを\nDB から取得", COLOR_PROCESS)
add_arrow(slide, 5, 5.0, 5, 5.4)

# PDF生成
add_box(slide, Inches(3.5), Inches(5.4), Inches(3), Inches(0.5), "PDFKit で PDF を生成\n（紫色ヘッダー）", COLOR_PROCESS)
add_arrow(slide, 5, 5.9, 5, 6.3)

# ダウンロード
add_box(slide, Inches(3.5), Inches(6.3), Inches(3), Inches(0.4), "ブラウザで PDF 表示", COLOR_PROCESS)
add_arrow(slide, 5, 6.7, 5, 7.0)

# 終了
add_oval(slide, Inches(3.75), Inches(7.0), Inches(1.5), Inches(0.5), "完了", COLOR_END)

# スライド 5: キャッシュメカニズム
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "グローバルキャッシュメカニズム")

# 初回ロード
add_box(slide, Inches(0.5), Inches(1.2), Inches(2.5), Inches(0.4), "projects.html\nロード", COLOR_PROCESS)
add_arrow(slide, 1.75, 1.6, 1.75, 2.0)

add_box(slide, Inches(0.5), Inches(2.0), Inches(2.5), Inches(0.4), "initializeGlobalCache()\n実行", COLOR_PROCESS)
add_arrow(slide, 1.75, 2.4, 1.75, 2.8)

add_box(slide, Inches(0.5), Inches(2.8), Inches(2.5), Inches(0.4), "GET /api/cache\n全データ取得", COLOR_PROCESS)
add_arrow(slide, 1.75, 3.2, 1.75, 3.6)

add_box(slide, Inches(0.5), Inches(3.6), Inches(2.5), Inches(0.5), "window.appCache\nに保存", COLOR_DECISION, RGBColor(0, 0, 0))
add_arrow(slide, 1.75, 4.1, 1.75, 4.5)

# 他ページ
add_box(slide, Inches(0.5), Inches(4.5), Inches(2.5), Inches(0.4), "他のページロード", COLOR_PROCESS)
add_arrow(slide, 1.75, 4.9, 1.75, 5.3)

add_box(slide, Inches(0.5), Inches(5.3), Inches(2.5), Inches(0.5), "window.appCache\nから読み込み", COLOR_PROCESS)
add_arrow(slide, 1.75, 5.8, 1.75, 6.2)

# テーブル描画
add_box(slide, Inches(0.5), Inches(6.2), Inches(2.5), Inches(0.4), "テーブル描画", COLOR_PROCESS)

# データ更新時
add_box(slide, Inches(3.5), Inches(1.2), Inches(2.5), Inches(0.4), "データ更新\n（POST/PUT/DELETE）", COLOR_PROCESS)
add_arrow(slide, 4.75, 1.6, 4.75, 2.0)

add_box(slide, Inches(3.5), Inches(2.0), Inches(2.5), Inches(0.4), "API 実行", COLOR_PROCESS)
add_arrow(slide, 4.75, 2.4, 4.75, 2.8)

add_box(slide, Inches(3.5), Inches(2.8), Inches(2.5), Inches(0.4), "キャッシュ再取得", COLOR_PROCESS)
add_arrow(slide, 4.75, 3.2, 4.75, 3.6)

add_box(slide, Inches(3.5), Inches(3.6), Inches(2.5), Inches(0.5), "window.appCache\n更新", COLOR_DECISION, RGBColor(0, 0, 0))
add_arrow(slide, 4.75, 4.1, 4.75, 4.5)

add_box(slide, Inches(3.5), Inches(4.5), Inches(2.5), Inches(0.4), "UI 再描画", COLOR_PROCESS)

# キャッシュボックス
add_box(slide, Inches(6.5), Inches(2.5), Inches(3), Inches(3.5),
        "window.appCache\n{\n  projects: [...]\n  vendors: [...]\n  categories: [...]\n  orders: [...]\n}",
        COLOR_PURPLE, RGBColor(255, 255, 255))

# スライド 6: データフロー全体
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "全体データフロー")

# ユーザー
add_oval(slide, Inches(4), Inches(1.0), Inches(2), Inches(0.5), "ユーザー操作", COLOR_START)
add_arrow(slide, 5, 1.5, 5, 1.9)

# window.appCache
add_box(slide, Inches(3.5), Inches(1.9), Inches(3), Inches(0.6), "window.appCache\n（グローバル）", RGBColor(156, 39, 176), RGBColor(255, 255, 255))
add_arrow(slide, 5, 2.5, 5, 2.9)

# Express Server
add_box(slide, Inches(3.5), Inches(2.9), Inches(3), Inches(0.6), "Express Server\n(Node.js)", COLOR_PROCESS)
add_arrow(slide, 5, 3.5, 5, 3.9)

# API Endpoints
add_box(slide, Inches(2.5), Inches(3.9), Inches(5), Inches(0.6),
        "POST/PUT/DELETE /api/projects/vendors/categories/orders | GET /api/cache | GET /api/invoice",
        RGBColor(255, 152, 0), RGBColor(255, 255, 255))
add_arrow(slide, 5, 4.5, 5, 4.9)

# SQLite Database
add_box(slide, Inches(3.5), Inches(4.9), Inches(3), Inches(0.6), "SQLite Database\n(better-sqlite3)", COLOR_END)
add_arrow(slide, 3.5, 5.5, 3.5, 5.9)

# テーブル
add_box(slide, Inches(1.5), Inches(5.9), Inches(1.5), Inches(0.4), "projects", COLOR_LIGHT, COLOR_DARK)
add_box(slide, Inches(3.2), Inches(5.9), Inches(1.5), Inches(0.4), "orders", COLOR_LIGHT, COLOR_DARK)
add_box(slide, Inches(4.9), Inches(5.9), Inches(1.5), Inches(0.4), "vendors", COLOR_LIGHT, COLOR_DARK)
add_box(slide, Inches(6.6), Inches(5.9), Inches(1.5), Inches(0.4), "categories", COLOR_LIGHT, COLOR_DARK)

# スライド 7: API エンドポイント
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "API エンドポイント & キャッシュ戦略")

y = 1.2
endpoints = [
    ("マスタ管理", [
        "GET /api/projects → 全案件取得",
        "POST /api/projects → 案件新規登録",
        "PUT /api/projects/:id → 案件更新",
        "DELETE /api/projects/:id → 案件削除"
    ]),
    ("注文・請求書", [
        "GET /api/orders → 全注文取得",
        "POST /api/orders → 注文新規登録",
        "GET /api/invoice/project/:id → 請求書PDF生成"
    ]),
    ("キャッシュ", [
        "GET /api/cache → 全マスタデータ一括取得",
        "返却: { projects, vendors, categories, orders }"
    ])
]

for category, items in endpoints:
    # カテゴリ名
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.25))
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = f"■ {category}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_PURPLE
    y += 0.35

    # エンドポイント
    for item in items:
        text_box = slide.shapes.add_textbox(Inches(1.0), Inches(y), Inches(8.5), Inches(0.25))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_DARK
        y += 0.35

    y += 0.15

# スライド 8: 請求書PDF レイアウト
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "請求書PDF のレイアウト")

# ヘッダー（紫色）
header = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(1), Inches(1.2), Inches(8), Inches(0.8)
)
header.fill.solid()
header.fill.fore_color.rgb = COLOR_PURPLE
header.line.width = Pt(0)

title_frame = header.text_frame
title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
p = title_frame.paragraphs[0]
p.text = "CONSTRUCT_PRO  請 求 書  2026年06月18日発行"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# 発注者情報
y = 2.1
info_items = [
    "【ご請求先】",
    "京都工務店 営業部",
    "電話: 075-123-4567 | メール: contact@kyoto-koumuten.jp",
    "住所: 京都府京都市中京区",
    "",
    "【案件名】",
    "京都御所南マンション改修工事"
]

for item in info_items:
    text_box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(0.25))
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = item
    if item.startswith("【"):
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_PURPLE
    else:
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_DARK
    y += 0.25

# テーブル示唆
y = 4.2
text_box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(0.25))
text_frame = text_box.text_frame
p = text_frame.paragraphs[0]
p.text = "【注文詳細】"
p.font.bold = True
p.font.size = Pt(10)
p.font.color.rgb = COLOR_PURPLE

# テーブルヘッダー
y += 0.35
headers = ["工事区分", "外注先", "見積額", "予定金額", "決定金額"]
x = 1.2
for header in headers:
    text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(1.3), Inches(0.25))
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = header
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # ヘッダーの背景
    cell_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x - 0.1), Inches(y - 0.05), Inches(1.3), Inches(0.25)
    )
    cell_bg.fill.solid()
    cell_bg.fill.fore_color.rgb = COLOR_PROCESS
    cell_bg.line.width = Pt(0)
    slide.shapes._spTree.remove(cell_bg._element)
    slide.shapes._spTree.insert(2, cell_bg._element)

    x += 1.3

# データ行
y += 0.35
data_rows = [
    ["塗装工事", "A会社", "¥2,000k", "¥1,950k", "¥1,800k"],
    ["電気工事", "B会社", "¥1,500k", "¥1,400k", "¥1,200k"],
    ["設備工事", "C会社", "¥1,200k", "¥1,150k", "¥1,000k"]
]

for row in data_rows:
    x = 1.2
    for cell in row:
        text_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(1.3), Inches(0.25))
        text_frame = text_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = cell
        p.font.size = Pt(8)
        p.font.color.rgb = COLOR_DARK
        p.alignment = PP_ALIGN.CENTER
        x += 1.3
    y += 0.3

# 金額計算
y += 0.2
text_box = slide.shapes.add_textbox(Inches(5.5), Inches(y), Inches(2.5), Inches(0.25))
text_frame = text_box.text_frame
p = text_frame.paragraphs[0]
p.text = "小計: ¥4,000,000"
p.font.size = Pt(9)
p.font.color.rgb = COLOR_DARK
p.alignment = PP_ALIGN.RIGHT

y += 0.3
text_box = slide.shapes.add_textbox(Inches(5.5), Inches(y), Inches(2.5), Inches(0.25))
text_frame = text_box.text_frame
p = text_frame.paragraphs[0]
p.text = "消費税（10%）: ¥400,000"
p.font.size = Pt(9)
p.font.color.rgb = COLOR_DARK
p.alignment = PP_ALIGN.RIGHT

y += 0.3
text_box = slide.shapes.add_textbox(Inches(5.5), Inches(y), Inches(2.5), Inches(0.35))
text_frame = text_box.text_frame
p = text_frame.paragraphs[0]
p.text = "合計: ¥4,400,000"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = COLOR_PURPLE
p.alignment = PP_ALIGN.RIGHT

# スライド 9: まとめ
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "CONSTRUCT_PRO 業務フロー")

# メインフロー
steps = [
    ("1. マスタ設定", Inches(0.5)),
    ("2. 案件登録", Inches(2.3)),
    ("3. 注文登録", Inches(4.1)),
    ("4. 請求書発行", Inches(5.9)),
    ("5. メール送信", Inches(7.7)),
]

for step, left in steps:
    add_box(slide, left, Inches(2.5), Inches(1.5), Inches(1.5), step, COLOR_PROCESS)

# 説明
y = 4.2
descriptions = [
    "案件、工事区分、外注先",
    "工事案件の基本情報",
    "工事の注文を登録",
    "発注元に請求書を生成",
    "PDF をメール送信"
]

for i, desc in enumerate(descriptions):
    left = Inches(0.5 + i * 1.8)
    text_box = slide.shapes.add_textbox(left, Inches(y), Inches(1.5), Inches(1.0))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = desc
    p.font.size = Pt(9)
    p.font.color.rgb = COLOR_DARK
    p.alignment = PP_ALIGN.CENTER

# ファイル保存
output_path = '/Users/hiromu.umeda/Desktop/construct-pro/CONSTRUCT_PRO_フローチャート.pptx'
prs.save(output_path)
print(f"✅ フローチャート付きPowerPoint ファイルを作成しました: {output_path}")
