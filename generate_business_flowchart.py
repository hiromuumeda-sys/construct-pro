#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

COLOR_PURPLE = RGBColor(139, 71, 137)
COLOR_PRIMARY = RGBColor(103, 80, 164)
COLOR_GREEN = RGBColor(76, 175, 80)
COLOR_BLUE = RGBColor(33, 150, 243)
COLOR_ORANGE = RGBColor(255, 152, 0)
COLOR_RED = RGBColor(244, 67, 54)
COLOR_DARK = RGBColor(49, 51, 56)
COLOR_LIGHT = RGBColor(245, 245, 245)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PURPLE

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

def add_title_bar(slide, title):
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.8)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_PRIMARY
    title_shape.line.color.rgb = COLOR_PRIMARY

    title_frame = title_shape.text_frame
    title_frame.clear()
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(5)

def add_process_box(slide, left, top, width, height, text, color, text_color=RGBColor(255, 255, 255)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)
    text_frame.margin_left = Inches(0.08)
    text_frame.margin_right = Inches(0.08)
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER

    return shape

def add_start_end(slide, left, top, width, height, text, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    return shape

def add_decision(slide, left, top, width, height, text, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    return shape

def add_arrow(slide, x1, y1, x2, y2, label=""):
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = COLOR_DARK
    connector.line.width = Pt(2)

    if label:
        mid_x = (x1 + x2) / 2
        mid_y = (y1 + y2) / 2
        text_box = slide.shapes.add_textbox(Inches(mid_x - 0.3), Inches(mid_y - 0.15), Inches(0.6), Inches(0.25))
        text_frame = text_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_DARK
        p.alignment = PP_ALIGN.CENTER

# ============== スライド作成 ==============

# スライド 1: タイトル
add_title_slide(prs, "CONSTRUCT_PRO", "業務フロー")

# スライド 2: 全体業務フロー
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "全体業務フロー")

# 1. 初期設定
add_process_box(slide, Inches(0.3), Inches(1.0), Inches(2.0), Inches(0.6), "① 初期設定\n（マスタ設定）", COLOR_GREEN)
add_arrow(slide, 1.3, 1.6, 1.3, 2.0)

# 2. 外注先登録
add_process_box(slide, Inches(0.3), Inches(2.0), Inches(2.0), Inches(0.6), "外注先を登録", COLOR_BLUE)
add_arrow(slide, 1.3, 2.6, 1.3, 3.0)

# 3. 工事区分登録
add_process_box(slide, Inches(0.3), Inches(3.0), Inches(2.0), Inches(0.6), "工事区分を登録", COLOR_BLUE)

# 矢印右へ
add_arrow(slide, 2.3, 3.3, 2.9, 3.3)

# 4. 案件登録
add_process_box(slide, Inches(2.9), Inches(1.0), Inches(2.0), Inches(0.6), "② 案件を登録", COLOR_BLUE)
add_arrow(slide, 3.9, 1.6, 3.9, 2.0)

# 説明
text_box = slide.shapes.add_textbox(Inches(2.9), Inches(2.0), Inches(2.0), Inches(1.0))
text_frame = text_box.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "• 案件名、発注元\n• 発注者情報\n• 金額・工期"
p.font.size = Pt(8)
p.font.color.rgb = COLOR_DARK

add_arrow(slide, 3.9, 3.0, 3.9, 3.4)

# 5. 注文登録
add_process_box(slide, Inches(2.9), Inches(3.4), Inches(2.0), Inches(0.6), "③ 注文を登録", COLOR_BLUE)
add_arrow(slide, 3.9, 4.0, 3.9, 4.4)

# 説明
text_box = slide.shapes.add_textbox(Inches(2.9), Inches(4.4), Inches(2.0), Inches(0.8))
text_frame = text_box.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "• 工事区分選択\n• 外注先選択\n• 金額入力"
p.font.size = Pt(8)
p.font.color.rgb = COLOR_DARK

add_arrow(slide, 3.9, 5.2, 3.9, 5.6)

# 6. 注文書送付
add_process_box(slide, Inches(2.9), Inches(5.6), Inches(2.0), Inches(0.6), "④ 注文書を\n発注先に送付", COLOR_ORANGE)
add_arrow(slide, 4.9, 5.9, 5.5, 5.9)

# 説明
text_box = slide.shapes.add_textbox(Inches(2.9), Inches(6.3), Inches(2.0), Inches(0.7))
text_frame = text_box.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "• 注文内容を PDF で作成\n• メール送信"
p.font.size = Pt(8)
p.font.color.rgb = COLOR_DARK

# 7. 請求書発行
add_process_box(slide, Inches(5.5), Inches(5.3), Inches(2.0), Inches(0.6), "⑤ 請求書を\n発行", COLOR_ORANGE)
add_arrow(slide, 6.5, 5.9, 6.5, 6.3)

# 説明
text_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(2.0), Inches(0.8))
text_frame = text_box.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "• 全注文を\n集計\n• PDF生成\n• 発注元に送付"
p.font.size = Pt(8)
p.font.color.rgb = COLOR_DARK

# 8. メール送信
add_process_box(slide, Inches(5.5), Inches(6.3), Inches(2.0), Inches(0.6), "⑥ 請求書を\nメール送信", COLOR_RED)
add_arrow(slide, 6.5, 6.9, 6.5, 7.0)

# 9. 支払い管理
add_process_box(slide, Inches(5.5), Inches(6.8), Inches(2.0), Inches(0.5), "⑦ 支払い完了", COLOR_RED)

# スライド 3: 初期設定フェーズ
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "① 初期設定フェーズ（マスタ設定）")

add_start_end(slide, Inches(3.75), Inches(1.0), Inches(1.5), Inches(0.5), "開始", COLOR_GREEN)
add_arrow(slide, 4.5, 1.5, 4.5, 2.0)

add_process_box(slide, Inches(3.5), Inches(2.0), Inches(2), Inches(0.6), "外注先マスタを\nセットアップ", COLOR_BLUE)
add_arrow(slide, 4.5, 2.6, 4.5, 3.1)

# 外注先情報
text_box = slide.shapes.add_textbox(Inches(5.8), Inches(1.8), Inches(3.7), Inches(1.2))
text_frame = text_box.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "外注先情報：\n会社名、部署、担当者\n電話、メール、住所"
p.font.size = Pt(10)
p.font.color.rgb = COLOR_DARK
p.font.bold = True

add_process_box(slide, Inches(3.5), Inches(3.1), Inches(2), Inches(0.6), "工事区分マスタを\nセットアップ", COLOR_BLUE)
add_arrow(slide, 4.5, 3.7, 4.5, 4.2)

# 工事区分情報
text_box = slide.shapes.add_textbox(Inches(5.8), Inches(2.9), Inches(3.7), Inches(1.2))
text_frame = text_box.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "工事区分：\n塗装工事、電気工事、\n設備工事など"
p.font.size = Pt(10)
p.font.color.rgb = COLOR_DARK
p.font.bold = True

add_process_box(slide, Inches(3.5), Inches(4.2), Inches(2), Inches(0.6), "マスタ設定完了", COLOR_GREEN)
add_arrow(slide, 4.5, 4.8, 4.5, 5.3)

add_start_end(slide, Inches(3.75), Inches(5.3), Inches(1.5), Inches(0.5), "完了", COLOR_GREEN)

# スライド 4: 案件登録フェーズ
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "② 案件登録フェーズ")

add_start_end(slide, Inches(3.75), Inches(1.0), Inches(1.5), Inches(0.5), "開始", COLOR_GREEN)
add_arrow(slide, 4.5, 1.5, 4.5, 2.0)

add_process_box(slide, Inches(3.5), Inches(2.0), Inches(2), Inches(0.6), "案件を新規登録", COLOR_BLUE)
add_arrow(slide, 4.5, 2.6, 4.5, 3.1)

# 案件情報
text_box = slide.shapes.add_textbox(Inches(5.8), Inches(1.8), Inches(3.7), Inches(1.5))
text_frame = text_box.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "案件情報：\n案件名、発注元\n発注者電話・メール・住所\n契約金額、工期"
p.font.size = Pt(10)
p.font.color.rgb = COLOR_DARK
p.font.bold = True

add_process_box(slide, Inches(3.5), Inches(3.1), Inches(2), Inches(0.6), "発注者情報を\n保存", COLOR_BLUE)
add_arrow(slide, 4.5, 3.7, 4.5, 4.2)

add_process_box(slide, Inches(3.5), Inches(4.2), Inches(2), Inches(0.6), "案件登録完了", COLOR_GREEN)
add_arrow(slide, 4.5, 4.8, 4.5, 5.3)

add_start_end(slide, Inches(3.75), Inches(5.3), Inches(1.5), Inches(0.5), "完了", COLOR_GREEN)

# スライド 5: 注文登録フェーズ
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "③ 注文登録フェーズ")

add_start_end(slide, Inches(4), Inches(1.0), Inches(1.5), Inches(0.5), "開始", COLOR_GREEN)
add_arrow(slide, 4.75, 1.5, 4.75, 2.0)

add_process_box(slide, Inches(3.75), Inches(2.0), Inches(2), Inches(0.6), "案件を選択", COLOR_BLUE)
add_arrow(slide, 4.75, 2.6, 4.75, 3.1)

add_process_box(slide, Inches(3.75), Inches(3.1), Inches(2), Inches(0.6), "注文管理画面\nを開く", COLOR_BLUE)
add_arrow(slide, 4.75, 3.7, 4.75, 4.2)

add_process_box(slide, Inches(3.75), Inches(4.2), Inches(2), Inches(0.6), "工事区分を\n選択", COLOR_BLUE)
add_arrow(slide, 4.75, 4.8, 4.75, 5.3)

# 右側：外注先選択
add_process_box(slide, Inches(6.0), Inches(4.2), Inches(2), Inches(0.6), "外注先を\n選択", COLOR_BLUE)
add_arrow(slide, 6.5, 3.7, 7.0, 4.2)

# 矢印合流
add_arrow(slide, 7.0, 4.8, 6.75, 5.3)
add_arrow(slide, 4.75, 5.3, 6.75, 5.3)

# 金額入力
add_process_box(slide, Inches(5.75), Inches(5.3), Inches(2), Inches(0.6), "見積額・決定額\nを入力", COLOR_BLUE)
add_arrow(slide, 6.75, 5.9, 6.75, 6.4)

add_process_box(slide, Inches(5.75), Inches(6.4), Inches(2), Inches(0.6), "注文を登録", COLOR_GREEN)
add_arrow(slide, 6.75, 7.0, 6.75, 7.0)

# スライド 6: 注文書送付フェーズ
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "④ 注文書を発注先に送付フェーズ")

add_start_end(slide, Inches(3.75), Inches(1.0), Inches(1.5), Inches(0.5), "開始", COLOR_GREEN)
add_arrow(slide, 4.5, 1.5, 4.5, 2.0)

add_process_box(slide, Inches(3.5), Inches(2.0), Inches(2), Inches(0.6), "登録した注文情報\nを確認", COLOR_BLUE)
add_arrow(slide, 4.5, 2.6, 4.5, 3.1)

add_process_box(slide, Inches(3.5), Inches(3.1), Inches(2), Inches(0.6), "工事区分・外注先\nごとに注文書\nを作成", COLOR_BLUE)
add_arrow(slide, 4.5, 3.7, 4.5, 4.2)

# 注文書情報
text_box = slide.shapes.add_textbox(Inches(5.8), Inches(2.8), Inches(3.7), Inches(2.0))
text_frame = text_box.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "注文書内容：\n• 案件名\n• 工事内容\n• 見積額・決定金額\n• 工期\n• 見積条件"
p.font.size = Pt(10)
p.font.color.rgb = COLOR_DARK
p.font.bold = True

add_process_box(slide, Inches(3.5), Inches(4.2), Inches(2), Inches(0.6), "注文書を\nPDF化", COLOR_ORANGE)
add_arrow(slide, 4.5, 4.8, 4.5, 5.3)

add_process_box(slide, Inches(3.5), Inches(5.3), Inches(2), Inches(0.6), "外注先に\nメール送信", COLOR_ORANGE)
add_arrow(slide, 4.5, 5.9, 4.5, 6.4)

add_start_end(slide, Inches(3.75), Inches(6.4), Inches(1.5), Inches(0.5), "完了", COLOR_GREEN)

# スライド 7: 請求書発行フェーズ
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "⑤ 請求書発行フェーズ")

add_start_end(slide, Inches(3.75), Inches(1.0), Inches(1.5), Inches(0.5), "開始", COLOR_GREEN)
add_arrow(slide, 4.5, 1.5, 4.5, 2.0)

add_process_box(slide, Inches(3.5), Inches(2.0), Inches(2), Inches(0.6), "案件の注文管理\n画面を開く", COLOR_BLUE)
add_arrow(slide, 4.5, 2.6, 4.5, 3.1)

add_process_box(slide, Inches(3.5), Inches(3.1), Inches(2), Inches(0.6), "「請求書を発行」\nをクリック", COLOR_BLUE)
add_arrow(slide, 4.5, 3.7, 4.5, 4.2)

add_process_box(slide, Inches(3.5), Inches(4.2), Inches(2), Inches(0.6), "請求書PDF\nを生成", COLOR_ORANGE)
add_arrow(slide, 4.5, 4.8, 4.5, 5.3)

# 請求書情報
text_box = slide.shapes.add_textbox(Inches(5.8), Inches(3.9), Inches(3.7), Inches(1.8))
text_frame = text_box.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "請求書内容：\n• 発注元情報\n• 工事ごとの注文\n• 決定額の合計\n• 消費税（10%）\n• 合計金額"
p.font.size = Pt(9)
p.font.color.rgb = COLOR_DARK
p.font.bold = True

add_process_box(slide, Inches(3.5), Inches(5.3), Inches(2), Inches(0.6), "PDF をダウンロード\nまたは印刷", COLOR_ORANGE)
add_arrow(slide, 4.5, 5.9, 4.5, 6.4)

add_start_end(slide, Inches(3.75), Inches(6.4), Inches(1.5), Inches(0.5), "完了", COLOR_GREEN)

# スライド 8: メール送信 & 支払い管理フェーズ
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "⑥⑦ 請求書メール送信 & 支払い管理フェーズ")

# メール送信
y = 1.2
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(4), Inches(0.3))
text_frame = text_box.text_frame
p = text_frame.paragraphs[0]
p.text = "【メール送信】（準備中）"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLOR_PURPLE

add_start_end(slide, Inches(1), Inches(1.6), Inches(1.3), Inches(0.4), "開始", COLOR_GREEN)
add_arrow(slide, 1.65, 2.0, 1.65, 2.4)

add_process_box(slide, Inches(0.8), Inches(2.4), Inches(1.7), Inches(0.5), "メールを\n作成", COLOR_BLUE)
add_arrow(slide, 1.65, 2.9, 1.65, 3.3)

add_process_box(slide, Inches(0.8), Inches(3.3), Inches(1.7), Inches(0.5), "請求書PDF\nを添付", COLOR_BLUE)
add_arrow(slide, 1.65, 3.8, 1.65, 4.2)

add_process_box(slide, Inches(0.8), Inches(4.2), Inches(1.7), Inches(0.5), "発注元に\n送信", COLOR_ORANGE)
add_arrow(slide, 1.65, 4.7, 1.65, 5.1)

add_start_end(slide, Inches(1), Inches(5.1), Inches(1.3), Inches(0.4), "完了", COLOR_GREEN)

# 支払い管理
y = 1.2
text_box = slide.shapes.add_textbox(Inches(5.5), Inches(y), Inches(4), Inches(0.3))
text_frame = text_box.text_frame
p = text_frame.paragraphs[0]
p.text = "【支払い管理】"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLOR_PURPLE

add_start_end(slide, Inches(6.1), Inches(1.6), Inches(1.3), Inches(0.4), "開始", COLOR_GREEN)
add_arrow(slide, 6.75, 2.0, 6.75, 2.4)

add_process_box(slide, Inches(5.8), Inches(2.4), Inches(1.9), Inches(0.5), "案件マスタ\nを開く", COLOR_BLUE)
add_arrow(slide, 6.75, 2.9, 6.75, 3.3)

add_process_box(slide, Inches(5.8), Inches(3.3), Inches(1.9), Inches(0.5), "支払い状況\nを確認", COLOR_BLUE)
add_arrow(slide, 6.75, 3.8, 6.75, 4.2)

add_decision(slide, Inches(5.95), Inches(4.2), Inches(1.6), Inches(0.6), "支払い済み？", COLOR_ORANGE)
add_arrow(slide, 6.75, 4.8, 6.75, 5.2)

add_process_box(slide, Inches(5.8), Inches(5.2), Inches(1.9), Inches(0.5), "「支払い済み」\nにチェック", COLOR_RED)
add_arrow(slide, 6.75, 5.7, 6.75, 6.1)

add_start_end(slide, Inches(6.1), Inches(6.1), Inches(1.3), Inches(0.4), "完了", COLOR_GREEN)

# スライド 9: データ参照フロー
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "各画面での実行内容")

# テーブルを作成
y = 1.0
content = [
    ["画面名", "実行内容", "データ参照"],
    ["案件マスタ", "• 案件登録・編集・削除\n• 請求書発行\n• 支払い管理", "外注先、工事区分"],
    ["工事区分マスタ", "• 工事区分の登録・編集\n• 削除", "（独立）"],
    ["外注先マスタ", "• 外注先の登録・編集\n• 削除", "（独立）"],
    ["注文管理", "• 注文登録・編集・削除\n• 注文書発行・送付\n• 請求書発行・送付", "案件、工事区分、外注先"]
]

x_positions = [0.5, 2.5, 5.5]
col_widths = [2, 3, 3]

for i, row in enumerate(content):
    for j, cell in enumerate(row):
        x = x_positions[j]
        w = col_widths[j]

        if i == 0:  # ヘッダー行
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y), Inches(w), Inches(0.35)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = COLOR_PRIMARY
            bg.line.width = Pt(1)
            bg.line.color.rgb = RGBColor(200, 200, 200)

            text_box = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y + 0.02), Inches(w - 0.1), Inches(0.31))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = cell
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
        else:  # データ行
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y), Inches(w), Inches(0.65)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = COLOR_LIGHT
            bg.line.width = Pt(1)
            bg.line.color.rgb = RGBColor(200, 200, 200)

            text_box = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y + 0.02), Inches(w - 0.1), Inches(0.61))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = cell
            p.font.size = Pt(9)
            p.font.color.rgb = COLOR_DARK

    if i == 0:
        y += 0.39
    else:
        y += 0.69

# スライド 10: システムの特徴
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "CONSTRUCT_PRO の特徴")

features = [
    ("1. グローバルキャッシュシステム", "初回読み込み時に全マスタデータを取得し、\n各ページで高速に参照可能", COLOR_BLUE),
    ("2. 注文書の自動生成・送付", "工事区分ごとに注文書を作成し、\n外注先にメール送信", COLOR_ORANGE),
    ("3. 一括請求書発行", "案件の全注文をまとめた請求書を\nPDF形式で自動生成", COLOR_ORANGE),
    ("4. 支払い管理", "案件ごとに支払い/未払いを管理し、\n可視化する仕組み", COLOR_RED),
]

y = 1.2
for title, desc, color in features:
    # タイトル
    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(8.6), Inches(0.25))
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = color

    # 説明
    text_box = slide.shapes.add_textbox(Inches(1.2), Inches(y + 0.3), Inches(8.1), Inches(0.6))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_DARK

    y += 1.0

# スライド 11: Railway アーキテクチャ
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "本番環境アーキテクチャ（Railway）")

# クライアント
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.3))
text_frame = text_box.text_frame
p = text_frame.paragraphs[0]
p.text = "【クライアント】"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLOR_PURPLE

add_process_box(slide, Inches(3.75), Inches(1.5), Inches(2.5), Inches(0.5), "ブラウザ\n(HTML/CSS/JS)", COLOR_BLUE)

add_arrow(slide, 5, 2.0, 5, 2.4)

# Railway App
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(9), Inches(0.3))
text_frame = text_box.text_frame
p = text_frame.paragraphs[0]
p.text = "【Railway クラウド】"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLOR_PURPLE

add_process_box(slide, Inches(3.5), Inches(2.8), Inches(3), Inches(0.6), "Node.js + Express\n(CONSTRUCT_PRO)", COLOR_BLUE)

# 上向き矢印
add_arrow(slide, 5, 3.4, 5, 3.8)

# 機能
y_func = 3.8
functions = [
    "API エンドポイント",
    "PDF 生成（PDFKit）",
    "メール送信（nodemailer）"
]
for func in functions:
    text_box = slide.shapes.add_textbox(Inches(5.8), Inches(y_func), Inches(3.5), Inches(0.25))
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = f"• {func}"
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_DARK
    y_func += 0.3

# 下向き矢印
add_arrow(slide, 5, 3.4, 5, 4.3)

# Database
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(0.3))
text_frame = text_box.text_frame
p = text_frame.paragraphs[0]
p.text = "【Railway Database】"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = COLOR_PURPLE

add_process_box(slide, Inches(2.5), Inches(4.7), Inches(2), Inches(0.6), "PostgreSQL\n（本番DB）", COLOR_GREEN)

add_arrow(slide, 3.5, 5.3, 3.5, 5.7)

# テーブル
y_db = 5.7
tables = ["projects", "orders", "vendors", "categories"]
for i, table in enumerate(tables):
    add_process_box(slide, Inches(1 + i * 1.8), Inches(y_db), Inches(1.5), Inches(0.4), table, COLOR_LIGHT, COLOR_DARK)

# 外部サービス
add_arrow(slide, 6.5, 5.3, 7.5, 5.5)
add_process_box(slide, Inches(7.2), Inches(4.7), Inches(2), Inches(0.6), "SendGrid\nまたは\nAWS SES", COLOR_ORANGE)

# S3 ストレージ
add_arrow(slide, 6.5, 3.4, 7.5, 3.6)
add_process_box(slide, Inches(7.2), Inches(2.8), Inches(2), Inches(0.6), "AWS S3\n（PDF・ファイル\n保存）", COLOR_ORANGE)

# スライド 12: 運用コスト内訳
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "本番環境ランニングコスト内訳（Railway）")

# コスト詳細
y = 1.1
costs = [
    ("1. Railway App", [
        "• CPU: 0.5 ~ 2 vCPU（スケール可能）",
        "• メモリ: 0.5GB ~ 2GB（スケール可能）",
        "• 従量課金: 約¥30～100/月（軽負荷時）",
        "  → ¥300～600/月（中負荷）"
    ]),
    ("2. Railway PostgreSQL", [
        "• ストレージ: 1GB（含まれる）",
        "• 追加容量: ¥2～5/GB",
        "• 従量課金: 約¥50～150/月（軽～中負荷）"
    ]),
    ("3. AWS S3（PDF・ファイル保存）", [
        "• ストレージ: 最初の1GB無料",
        "• 従量課金: ¥0.023/GB（東京リージョン）",
        "• 請求額: 約¥0～50/月（軽～中負荷）"
    ]),
    ("4. メール送信サービス", [
        "• SendGrid: 月100件無料、以降¥2～10/月",
        "• AWS SES: 月62,000件無料、以降¥0.1/件",
        "• 推定: ¥0～100/月（本番想定）"
    ]),
    ("5. DB バックアップ（デイリー）", [
        "• Railway バックアップ: 約¥700～1,400/月",
        "• 保持期間: 直近30日分保存",
        "• 自動復旧機能付き"
    ]),
    ("6. ドメイン・SSL", [
        "• Railway カスタムドメイン: 無料",
        "• SSL証明書: 自動更新（無料）",
        "• 外部ドメイン: ¥500～1,500/年"
    ])
]

for category, items in costs:
    # カテゴリ
    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(8.6), Inches(0.28))
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = category
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    y += 0.32

    # 詳細
    for item in items:
        text_box = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(8.1), Inches(0.25))
        text_frame = text_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_DARK
        y += 0.26

    y += 0.1

# 合計
y += 0.2
add_process_box(slide, Inches(1.5), Inches(y), Inches(7), Inches(0.5), "想定月額コスト: ¥780～2,350/月（デイリーバックアップ含む）", COLOR_PURPLE)

# スライド 13: データベース移行ガイド
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "SQLite → PostgreSQL 移行ガイド")

steps = [
    ("ステップ 1: Railway 準備", [
        "Railway にプロジェクト作成",
        "PostgreSQL アドオンを追加",
        "接続情報（HOST, PORT, USER, PASS）を取得"
    ]),
    ("ステップ 2: コード修正", [
        "better-sqlite3 → node-postgres (pg) に切り替え",
        "環境変数から DB 接続情報を読み込み",
        "スキーマ（テーブル定義）を実行"
    ]),
    ("ステップ 3: データ移行", [
        "既存 SQLite データをエクスポート",
        "PostgreSQL にインポート",
        "データ整合性を確認"
    ]),
    ("ステップ 4: デプロイ", [
        "GitHub に code push",
        "Railway が自動ビルド・デプロイ",
        "本番環境で動作確認"
    ])
]

y = 1.1
for step_name, details in steps:
    # ステップ名
    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(8.6), Inches(0.28))
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = step_name
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    y += 0.32

    # 詳細
    for detail in details:
        text_box = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(8.1), Inches(0.24))
        text_frame = text_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = detail
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_DARK
        y += 0.25

    y += 0.08

# スライド 14: 本番チェックリスト
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "本番デプロイ前チェックリスト")

checklist = [
    "□ Railway アカウント作成・プロジェクト設定",
    "□ PostgreSQL アドオン追加・初期化",
    "□ 環境変数設定（DB_HOST, DB_USER, DB_PASS など）",
    "□ GitHub リポジトリ接続",
    "□ .env.production ファイル準備",
    "□ node_modules, .git などの.gitignore 確認",
    "□ SMTP サービス設定（メール送信機能用）",
    "□ ドメイン設定（カスタムドメイン or railway.app 利用）",
    "□ SSL 証明書（自動更新に任せる）",
    "□ ローカルテスト完了（全機能動作確認）",
    "□ セキュリティ確認（SQL インジェクション対策など）",
    "□ バックアップ戦略決定（Railway バックアップ機能利用）",
]

y = 1.0
for item in checklist:
    text_box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(0.25))
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = item
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_DARK
    y += 0.3

# スライド 15: 本番運用の注意点
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_bar(slide, "本番運用上の注意点")

cautions = [
    ("セキュリティ", [
        "• 環境変数に機密情報（DB パス、API キー）を保存",
        "• HTTPS 必須（Railway が自動対応）",
        "• SQL インジェクション対策：パラメータ化クエリを使用"
    ]),
    ("バックアップ & 災害復旧", [
        "• Railway PostgreSQL デイリーバックアップ（自動）",
        "• 保持期間: 直近30日分（毎日実行）",
        "• バックアップから1時間以内に復旧可能",
        "• 月額コスト: ¥700～1,400（専用アドオン利用時）"
    ]),
    ("モニタリング", [
        "• Railway のログ機能で エラー監視",
        "• CPU・メモリ使用率を定期確認",
        "• メモリ不足時はスケール検討"
    ]),
    ("パフォーマンス", [
        "• データベースインデックス設定",
        "• クエリ最適化（不要な JOIN を削減）",
        "• キャッシング戦略（グローバルキャッシュを活用）"
    ])
]

y = 1.0
for category, items in cautions:
    # カテゴリ
    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(8.6), Inches(0.25))
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = f"■ {category}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_PURPLE
    y += 0.3

    # 詳細
    for item in items:
        text_box = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(8.1), Inches(0.24))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_DARK
        y += 0.25

    y += 0.15

# ファイル保存
output_path = '/Users/hiromu.umeda/Desktop/construct-pro/CONSTRUCT_PRO_業務フロー.pptx'
prs.save(output_path)
print(f"✅ 業務フロー版PowerPoint ファイルを作成しました: {output_path}")
