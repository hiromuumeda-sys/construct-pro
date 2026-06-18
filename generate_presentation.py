#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# プレゼンテーション作成
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# カラースキーム
COLOR_PURPLE = RGBColor(139, 71, 137)  # #8B4789
COLOR_PRIMARY = RGBColor(103, 80, 164)  # #6750A4
COLOR_DARK = RGBColor(49, 51, 56)
COLOR_LIGHT = RGBColor(245, 245, 245)

def add_title_slide(prs, title, subtitle):
    """タイトルスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PURPLE

    # タイトル
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # サブタイトル
    left = Inches(0.5)
    top = Inches(4.2)
    width = Inches(9)
    height = Inches(1)
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items):
    """コンテンツスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # タイトルバー
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_PRIMARY
    title_shape.line.color.rgb = COLOR_PRIMARY

    # タイトルテキスト
    title_frame = title_shape.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(8)

    # コンテンツ
    left = Inches(0.5)
    top = Inches(1.2)
    for i, item in enumerate(content_items):
        if isinstance(item, dict):
            # 見出し
            if item.get('type') == 'heading':
                text_box = slide.shapes.add_textbox(left, top, Inches(9), Inches(0.4))
                text_frame = text_box.text_frame
                p = text_frame.paragraphs[0]
                p.text = item['text']
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = COLOR_PURPLE
                top += Inches(0.5)
            # テキスト
            elif item.get('type') == 'text':
                indent = item.get('indent', 0) * Inches(0.3)
                text_box = slide.shapes.add_textbox(left + indent, top, Inches(9) - indent, Inches(0.3))
                text_frame = text_box.text_frame
                p = text_frame.paragraphs[0]
                p.text = item['text']
                p.font.size = Pt(14)
                p.font.color.rgb = COLOR_DARK
                top += Inches(0.35)
        else:
            # 単純なテキスト
            text_box = slide.shapes.add_textbox(left, top, Inches(9), Inches(0.3))
            text_frame = text_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = COLOR_DARK
            top += Inches(0.4)

def add_flowchart_slide(prs, title, flow_items):
    """フロー図スライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # タイトルバー
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_PRIMARY
    title_shape.line.color.rgb = COLOR_PRIMARY

    title_frame = title_shape.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # フロー図
    y = 1.5
    for item in flow_items:
        # ボックス
        box = slide.shapes.add_shape(1, Inches(1.5), Inches(y), Inches(7), Inches(0.6))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_LIGHT
        box.line.color.rgb = COLOR_PRIMARY
        box.line.width = Pt(2)

        # テキスト
        text_frame = box.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_DARK
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(6)
        p.space_after = Pt(6)

        y += Inches(0.9)

        # 矢印（最後でなければ）
        if item != flow_items[-1]:
            arrow = slide.shapes.add_connector(1, Inches(5), Inches(y - 0.3), Inches(5), Inches(y))
            arrow.line.color.rgb = COLOR_PRIMARY
            arrow.line.width = Pt(2)

# ============== スライド作成開始 ==============

# スライド 1: タイトル
add_title_slide(prs, "CONSTRUCT_PRO", "システムフロー")

# スライド 2: システム全体構成
add_content_slide(prs, "システム全体構成", [
    {'type': 'heading', 'text': '3つの主要フェーズ'},
    {'type': 'text', 'text': 'マスタ設定 - 案件、工事区分、外注先の初期設定', 'indent': 0},
    {'type': 'text', 'text': '注文管理 - 案件ごとの注文登録・編集', 'indent': 0},
    {'type': 'text', 'text': '請求・支払い - 請求書発行と支払い管理', 'indent': 0},
    {'type': 'heading', 'text': 'マスタデータ'},
    {'type': 'text', 'text': '案件マスタ - 工事案件の基本情報', 'indent': 0},
    {'type': 'text', 'text': '工事区分マスタ - 工事分類', 'indent': 0},
    {'type': 'text', 'text': '外注先マスタ - 外注先企業情報', 'indent': 0},
])

# スライド 3: 新規案件登録フロー
add_flowchart_slide(prs, "新規案件登録フロー", [
    "① モーダルで案件名、発注元、発注者情報を入力",
    "② 「登録」ボタンをクリック",
    "③ サーバーで案件データを検証、DBに保存",
    "④ キャッシュを再取得（GET /api/cache）",
    "⑤ テーブルに新規案件が表示される"
])

# スライド 4: データフロー
add_content_slide(prs, "データフロー", [
    {'type': 'heading', 'text': 'フロント → バックエンド'},
    {'type': 'text', 'text': 'POST /api/projects - 案件新規登録', 'indent': 1},
    {'type': 'text', 'text': 'PUT /api/projects/:id - 案件更新', 'indent': 1},
    {'type': 'text', 'text': 'DELETE /api/projects/:id - 案件削除', 'indent': 1},
    {'type': 'heading', 'text': 'バックエンド → フロント'},
    {'type': 'text', 'text': 'GET /api/cache - 全マスタデータ取得', 'indent': 1},
    {'type': 'text', 'text': '返却形式: { projects: [...], vendors: [...], ... }', 'indent': 1},
    {'type': 'heading', 'text': 'グローバルキャッシュ'},
    {'type': 'text', 'text': 'window.appCache に全データを保存', 'indent': 1},
    {'type': 'text', 'text': '他ページは window.appCache から参照', 'indent': 1},
])

# スライド 5: 各画面説明（案件マスタ）
add_content_slide(prs, "案件マスタ画面", [
    {'type': 'heading', 'text': '主要機能'},
    {'type': 'text', 'text': '案件の新規登録・編集・削除', 'indent': 1},
    {'type': 'text', 'text': '支払い済み/未払いの管理', 'indent': 1},
    {'type': 'text', 'text': '請求書発行', 'indent': 1},
    {'type': 'text', 'text': 'キーワード検索・ステータスフィルタ', 'indent': 1},
    {'type': 'heading', 'text': '登録項目'},
    {'type': 'text', 'text': '案件名、発注元（企業名）', 'indent': 1},
    {'type': 'text', 'text': '発注者: 電話、メール、住所', 'indent': 1},
    {'type': 'text', 'text': '契約金額、工期（開始・終了）', 'indent': 1},
])

# スライド 6: 工事区分マスタ
add_content_slide(prs, "工事区分マスタ画面", [
    {'type': 'heading', 'text': '目的'},
    {'type': 'text', 'text': '工事の分類を管理', 'indent': 1},
    {'type': 'heading', 'text': '主要機能'},
    {'type': 'text', 'text': '区分コード、区分名の登録', 'indent': 1},
    {'type': 'text', 'text': '表示順序の管理', 'indent': 1},
    {'type': 'text', 'text': '備考情報', 'indent': 1},
    {'type': 'heading', 'text': '用途'},
    {'type': 'text', 'text': '注文登録時に工事区分を選択', 'indent': 1},
])

# スライド 7: 外注先マスタ
add_content_slide(prs, "外注先マスタ画面", [
    {'type': 'heading', 'text': '目的'},
    {'type': 'text', 'text': '外注先企業の情報を管理', 'indent': 1},
    {'type': 'heading', 'text': '登録情報'},
    {'type': 'text', 'text': '会社名、部署、担当者', 'indent': 1},
    {'type': 'text', 'text': '電話、メール、住所', 'indent': 1},
    {'type': 'text', 'text': 'エリア（地域分類）', 'indent': 1},
    {'type': 'heading', 'text': '用途'},
    {'type': 'text', 'text': '注文登録時に外注先を選択', 'indent': 1},
])

# スライド 8: 注文管理フロー
add_content_slide(prs, "注文管理フロー", [
    {'type': 'heading', 'text': 'ステップ1: 案件選択'},
    {'type': 'text', 'text': '案件マスタから案件を選択', 'indent': 1},
    {'type': 'text', 'text': '「注文管理」ボタンをクリック', 'indent': 1},
    {'type': 'heading', 'text': 'ステップ2: 注文登録'},
    {'type': 'text', 'text': '工事区分、外注先、金額を入力', 'indent': 1},
    {'type': 'text', 'text': '複数の注文を登録可能', 'indent': 1},
    {'type': 'heading', 'text': 'ステップ3: 請求書発行'},
    {'type': 'text', 'text': '案件の全注文をまとめた請求書を生成', 'indent': 1},
])

# スライド 9: 請求書発行プロセス
add_flowchart_slide(prs, "請求書発行プロセス", [
    "① 注文管理画面で「請求書を発行」ボタンをクリック",
    "② サーバーが案件と関連注文をDB から取得",
    "③ PDFKit で請求書を生成（紫色ヘッダー）",
    "④ テーブルに注文情報を整形",
    "⑤ 金額計算（小計 + 消費税10%）",
    "⑥ ブラウザで PDF をダウンロード / 表示"
])

# スライド 10: 請求書PDF のレイアウト
add_content_slide(prs, "請求書PDF のレイアウト", [
    {'type': 'heading', 'text': 'ヘッダー'},
    {'type': 'text', 'text': 'CONSTRUCT_PRO（紫色背景 #8B4789）', 'indent': 1},
    {'type': 'text', 'text': '発行日付表示', 'indent': 1},
    {'type': 'heading', 'text': '発注者情報'},
    {'type': 'text', 'text': '企業名、部署、担当者', 'indent': 1},
    {'type': 'text', 'text': '電話、メール、住所', 'indent': 1},
    {'type': 'heading', 'text': '注文詳細'},
    {'type': 'text', 'text': 'テーブル形式で注文情報を表示', 'indent': 1},
    {'type': 'heading', 'text': '金額計算'},
    {'type': 'text', 'text': '小計、消費税（10%）、合計', 'indent': 1},
])

# スライド 11: メール送信機能（予定）
add_content_slide(prs, "メール送信機能（準備中）", [
    {'type': 'heading', 'text': '概要'},
    {'type': 'text', 'text': '請求書をメールで発注元に送信', 'indent': 1},
    {'type': 'heading', 'text': 'エンドポイント'},
    {'type': 'text', 'text': 'POST /api/email/invoice', 'indent': 1},
    {'type': 'heading', 'text': 'リクエストボディ'},
    {'type': 'text', 'text': 'projectId, to, subject, body', 'indent': 1},
    {'type': 'heading', 'text': '実装技術'},
    {'type': 'text', 'text': 'nodemailer（SMTP 経由でメール送信）', 'indent': 1},
])

# スライド 12: キャッシュメカニズム
add_flowchart_slide(prs, "グローバルキャッシュメカニズム", [
    "① projects.html 読み込み時に initializeGlobalCache() 実行",
    "② GET /api/cache で全データ取得",
    "③ window.appCache に保存",
    "④ 他のページでは window.appCache から参照",
    "⑤ データ更新時は initializeGlobalCache() を再実行",
    "⑥ キャッシュを最新状態に保つ"
])

# スライド 13: エラー処理
add_content_slide(prs, "エラー処理とバリデーション", [
    {'type': 'heading', 'text': 'クライアント側'},
    {'type': 'text', 'text': '必須項目の空白チェック', 'indent': 1},
    {'type': 'text', 'text': '日付形式の検証', 'indent': 1},
    {'type': 'heading', 'text': 'サーバー側'},
    {'type': 'text', 'text': '入力データの検証', 'indent': 1},
    {'type': 'text', 'text': 'データベーストランザクション処理', 'indent': 1},
    {'type': 'heading', 'text': 'レスポンス形式'},
    {'type': 'text', 'text': '{ error: "error_type", message: "説明" }', 'indent': 1},
])

# スライド 14: 技術スタック
add_content_slide(prs, "技術スタック", [
    {'type': 'heading', 'text': 'フロントエンド'},
    {'type': 'text', 'text': 'HTML5, Tailwind CSS', 'indent': 1},
    {'type': 'heading', 'text': 'バックエンド'},
    {'type': 'text', 'text': 'Node.js + Express（ポート 3000）', 'indent': 1},
    {'type': 'heading', 'text': 'データベース'},
    {'type': 'text', 'text': 'SQLite（better-sqlite3）', 'indent': 1},
    {'type': 'heading', 'text': 'PDF 生成'},
    {'type': 'text', 'text': 'PDFKit', 'indent': 1},
    {'type': 'heading', 'text': 'メール'},
    {'type': 'text', 'text': 'nodemailer（SMTP）', 'indent': 1},
])

# スライド 15: 業務フロー例
add_flowchart_slide(prs, "典型的な業務フロー", [
    "1. 案件マスタで新規案件登録（京都御所南マンション改修工事）",
    "2. 案件の注文管理画面を開く",
    "3. 複数の注文を登録（塗装工事、電気工事、設備工事など）",
    "4. 「請求書を発行」で PDF 生成",
    "5. メールで発注元に送信（準備中）",
    "6. 支払い完了後、支払い済みにチェック"
])

# スライド 16: まとめ
add_content_slide(prs, "まとめ", [
    {'type': 'heading', 'text': 'CONSTRUCT_PRO の流れ'},
    {'type': 'text', 'text': 'マスタ設定（案件、工事区分、外注先）', 'indent': 0},
    {'type': 'text', 'text': '案件登録', 'indent': 0},
    {'type': 'text', 'text': '注文登録（複数可）', 'indent': 0},
    {'type': 'text', 'text': '請求書発行（まとめて生成）', 'indent': 0},
    {'type': 'text', 'text': 'メール送信（準備中）', 'indent': 0},
    {'type': 'text', 'text': '支払い管理（状況記録）', 'indent': 0},
    {'type': 'text', 'text': '完了', 'indent': 0},
])

# ファイル保存
output_path = '/Users/hiromu.umeda/Desktop/construct-pro/CONSTRUCT_PRO_システムフロー.pptx'
prs.save(output_path)
print(f"✅ PowerPoint ファイルを作成しました: {output_path}")
